import argparse
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF
import os
from typing import List

# Function to convert PDF to PPTX
def pdf_to_pptx(pdf_path: str, output_path: str, dpi: int) -> None:
    """
    Convert a PDF file into a PowerPoint presentation.

    Args:
        pdf_path (str): Path to the input PDF file.
        output_path (str): Path to save the output PPTX file.
        dpi (int): Resolution for rendering PDF pages to images.
    """
    print("Reading PDF file...")
    doc = fitz.open(pdf_path)  # Open the PDF document

    # Create a PowerPoint presentation
    print("Creating PowerPoint presentation...")
    presentation = Presentation()

    for i, page in enumerate(doc):
        # Render the page to a pixmap (image)
        pix = page.get_pixmap(dpi=dpi)
        image_path = f"temp_page_{i + 1}.png"
        pix.save(image_path)  # Save the image as a PNG file

        # Add a new slide to the presentation
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        # Add the image to the slide
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

        # Remove the temporary image file
        os.remove(image_path)

    # Save the presentation
    print("Saving PowerPoint presentation...")
    presentation.save(output_path)
    print(f"Presentation saved to {output_path}")

# Main function
def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convert a PDF file into a PowerPoint presentation.",
        epilog="Example usage: python script_name.py --pdf-path input.pdf --output-path output.pptx --dpi 300"
    )
    parser.add_argument("--pdf", type=str, required=True, help="Path to the input PDF file (e.g., /path/to/input.pdf).")
    parser.add_argument("--output", type=str, required=True, help="Path to save the output PPTX file (e.g., /path/to/output.pptx).")
    parser.add_argument("--dpi", type=int, default=300, help="Resolution for rendering PDF pages to images (default: 300).")

    args = parser.parse_args()

    if not os.path.exists(args.pdf):
        print(f"PDF file '{args.pdf}' not found. Please check the path.")
    else:
        pdf_to_pptx(args.pdf, args.output, args.dpi)


if __name__ == "__main__":
    main()
